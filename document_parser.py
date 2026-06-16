import os
import pypdf
import docx
import pptx

# Extensions we support
TEXT_EXTENSIONS = {
    '.txt', '.py', '.md', '.json', '.xml', '.csv', '.log', 
    '.js', '.ts', '.html', '.css', '.ini', '.cfg', '.yaml', 
    '.yml', '.sh', '.bat', '.ps1', '.c', '.cpp', '.h', '.java'
}
OFFICE_EXTENSIONS = {'.pdf', '.docx', '.pptx'}
ALL_SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS.union(OFFICE_EXTENSIONS)

def is_supported_file(filename):
    _, ext = os.path.splitext(filename.lower())
    return ext in ALL_SUPPORTED_EXTENSIONS

def extract_text_from_pdf(filepath):
    text_lines = []
    try:
        reader = pypdf.PdfReader(filepath)
        for page_no, page in enumerate(reader.pages, 1):
            page_text = page.extract_text()
            if page_text:
                for line in page_text.splitlines():
                    clean_line = line.strip()
                    if clean_line:
                        text_lines.append((clean_line, f"Page {page_no}"))
    except Exception as e:
        print(f"Error reading PDF {filepath}: {e}")
    return text_lines

def extract_text_from_docx(filepath):
    text_lines = []
    try:
        doc = docx.Document(filepath)
        # Parse paragraphs
        for para_no, para in enumerate(doc.paragraphs, 1):
            clean_text = para.text.strip()
            if clean_text:
                text_lines.append((clean_text, f"Paragraph {para_no}"))
                
        # Parse tables
        for table_no, table in enumerate(doc.tables, 1):
            for row_no, row in enumerate(table.rows, 1):
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_lines.append((row_text, f"Table {table_no} Row {row_no}"))
    except Exception as e:
        print(f"Error reading Word document {filepath}: {e}")
    return text_lines

def extract_text_from_pptx(filepath):
    text_lines = []
    try:
        prs = pptx.Presentation(filepath)
        for slide_no, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    for line in shape.text.splitlines():
                        clean_line = line.strip()
                        if clean_line:
                            text_lines.append((clean_line, f"Slide {slide_no}"))
    except Exception as e:
        print(f"Error reading PPTX {filepath}: {e}")
    return text_lines

def extract_text_from_file(filepath):
    """
    Extracts text from a file. 
    Returns list of tuples: [(line_content, reference_label)]
    Where reference_label can be "Line 4", "Page 2", "Slide 1", etc.
    """
    _, ext = os.path.splitext(filepath.lower())
    
    if ext in TEXT_EXTENSIONS:
        lines = []
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                for line_no, line in enumerate(f, 1):
                    clean_line = line.strip()
                    if clean_line:
                        lines.append((clean_line, f"Line {line_no}"))
        except Exception:
            pass
        return lines
        
    elif ext == '.pdf':
        return extract_text_from_pdf(filepath)
        
    elif ext == '.docx':
        return extract_text_from_docx(filepath)
        
    elif ext == '.pptx':
        return extract_text_from_pptx(filepath)
        
    return []
